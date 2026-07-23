"""Dedicated, allowlisted adapters for visible Microsoft desktop automation."""
from __future__ import annotations

import dataclasses
import json
import re
import time
from pathlib import Path

from config import Config
from core.live_task import TaskCancelled


@dataclasses.dataclass
class DesktopActionResult:
    status: str
    application: str
    action: str
    message: str
    file_path: str = ""
    application_left_open: bool = True
    error: str = ""
    recovery_attempt: str = ""


def descriptive_filename(value, extension):
    cleaned = re.sub(r"[<>:\"/\\|?*]+", " ", str(value or "")).strip()
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" .")
    cleaned = cleaned[:100] or "JARVIS Created File"
    suffix = extension if str(extension).startswith(".") else f".{extension}"
    if cleaned.lower().endswith(suffix.lower()):
        return cleaned
    return cleaned + suffix


def created_files_folder(kind):
    folder = Config.CREATED_FILES_DIR / str(kind)
    folder.mkdir(parents=True, exist_ok=True)
    return folder


def unique_path(path):
    candidate = Path(path)
    if not candidate.exists():
        return candidate
    for index in range(2, 1000):
        alternative = candidate.with_name(f"{candidate.stem} ({index}){candidate.suffix}")
        if not alternative.exists():
            return alternative
    raise RuntimeError(f"No unused filename is available for {candidate.name}")


class DesktopActionLogger:
    def write(self, *, command="", intent="", application="", action="",
              result="", error="", file_path="", recovery_attempt="",
              approval_status="not_required"):
        from core.action_manager import ActionManager
        entry = {
            "timestamp": time.time(),
            "voice_command": ActionManager._redact_text(command),
            "detected_intent": intent,
            "application": application,
            "action": action,
            "result": ActionManager._redact_text(result),
            "error": ActionManager._redact_text(error),
            "file_path": ActionManager._redact_text(file_path),
            "recovery_attempt": ActionManager._redact_text(recovery_attempt),
            "approval_status": approval_status,
        }
        Config.DESKTOP_ACTION_LOG_FILE.parent.mkdir(parents=True, exist_ok=True)
        with Config.DESKTOP_ACTION_LOG_FILE.open("a", encoding="utf-8") as stream:
            stream.write(json.dumps(entry, ensure_ascii=False) + "\n")


class DesktopApplicationAdapter:
    application = ""

    def __init__(self, ctx):
        self.ctx = ctx

    def open(self):
        controller = getattr(self.ctx, "windows_controller", None)
        if controller is None:
            from core.windows_controller import WindowsController
            controller = WindowsController(self.ctx)
        return controller.open_application(self.application)


class WordAdapter(DesktopApplicationAdapter):
    application = "Microsoft Word"

    STRUCTURES = {
        "letter": ("Subject", "Greeting", "Request", "Supporting Details", "Closing"),
        "proposal": ("Executive Summary", "Problem", "Proposed Solution", "Objectives", "Scope", "Deliverables", "Timeline", "Risks", "Success Measures", "Approval Request"),
        "research": ("Executive Summary", "Introduction", "Background", "Method", "Findings", "Discussion", "Recommendations", "References"),
        "report": ("Executive Summary", "Background", "Key Findings", "Analysis", "Recommendations", "Conclusion"),
    }

    def _service(self, create=False):
        service = self.ctx.state.get("word_service")
        if service is None and create:
            from skills.office_service import WordService
            service = WordService()
            service.open(visible=True)
            self.ctx.state["word_service"] = service
        return service

    def open_word(self):
        service = self._service(create=True)
        return service.open(visible=True) if service.app is None else True

    def create_blank_document(self):
        return self._service(create=True).new_document()

    def insert_text(self, text):
        return self._service(create=True).insert_text(text)

    def apply_heading(self, text, level=1):
        return self._service(create=True).apply_heading(text, level)

    def apply_paragraph_style(self, style_name="Normal"):
        return self._service(create=True).apply_paragraph_style(style_name)

    def insert_table(self, data):
        return self._service(create=True).insert_table(0, data=data)

    def insert_page_break(self):
        return self._service(create=True).insert_page_break()

    def add_page_numbers(self):
        return self._service(create=True).add_page_numbers()

    def save_docx(self, path):
        return self._service(create=True).save(path)

    def export_pdf(self, path):
        return self._service(create=True).export_pdf(path)

    def close_document(self, save=True):
        service = self._service()
        return service.close_document(save) if service is not None else False

    def _sections(self, document_type, topic):
        low = str(document_type).lower()
        structure = next((value for key, value in self.STRUCTURES.items() if key in low), None)
        structure = structure or ("Overview", "Details", "Recommendations", "Conclusion")
        if getattr(self.ctx.llm, "available", False):
            prompt = (
                f"Draft a professional {document_type} about {topic or document_type}. "
                "Use markdown headings and complete paragraphs. Do not include secrets."
            )
            drafted = self.ctx.llm.quick(prompt, max_tokens=2600)
            if drafted:
                sections = []
                heading = "Overview"
                body = []
                for line in drafted.splitlines():
                    if line.lstrip().startswith("#"):
                        if body:
                            sections.append((heading, "\n".join(body).strip()))
                        heading = line.lstrip("# ").strip() or heading
                        body = []
                    elif line.strip():
                        body.append(line.strip())
                if body:
                    sections.append((heading, "\n".join(body).strip()))
                if sections:
                    return sections
        subject = topic or document_type.title()
        return [
            (heading, f"{heading} for {subject}. This section is ready for review and further detail.")
            for heading in structure
        ]

    def create_document(self, document_type="document", topic="", mode="instant",
                        save_after_completion=True):
        from skills.office_service import WordService
        service = WordService()
        title = (topic or document_type).strip().title() or "JARVIS Document"
        filename = descriptive_filename(f"{title} {document_type.title()}", ".docx")
        path = unique_path(created_files_folder("Word") / filename)
        task = getattr(self.ctx, "live_task", None)
        try:
            if task is not None:
                task.start(f"word-{int(time.time())}", f"Create {document_type}", application=self.application, mode=mode.upper())
            service.open(visible=True)
            document = service.new_document()
            service.insert_heading(title, level=1, doc=document)
            sections = self._sections(document_type, topic)
            total = max(1, len(sections))
            for index, (heading, body) in enumerate(sections, 1):
                if task is not None:
                    task.checkpoint()
                    task.update(step=f"Writing {heading}", progress=int(index * 90 / total))
                service.insert_heading(heading, level=1, doc=document)
                service.type_visibly(body, doc=document)
                if mode in {"visible", "structured"}:
                    time.sleep(max(0, Config.LIVE_TYPING_DELAY_MS) / 1000.0)
            if save_after_completion:
                service.save(path, doc=document)
            if save_after_completion and not path.exists():
                raise RuntimeError("Word reported success but the DOCX file was not found")
            entry = self.ctx.registry.register(
                "document", path.name if save_after_completion else title,
                window_title=path.stem if save_after_completion else "Word",
                pid=service.process_id, hwnd=service.window_handle,
                closer=lambda svc=service: svc.close(save=True),
                extra={"path": str(path), "application": self.application,
                       "unsaved": "false" if save_after_completion else "true"},
            )
            self.ctx.state["word_service"] = service
            self.ctx.state["active_office_entry"] = entry["id"]
            if task is not None:
                task.complete()
            return DesktopActionResult(
                "success", self.application, "create_document",
                f"Created and verified {path.name}. Word remains open for review.",
                str(path), True,
            )
        except TaskCancelled:
            try:
                service.close(save=False)
            except Exception:
                pass
            return DesktopActionResult(
                "cancelled", self.application, "create_document",
                "Word document creation was cancelled.",
            )
        except Exception as exc:
            if task is not None:
                task.fail(exc)
            return DesktopActionResult("failed", self.application, "create_document", "Word document creation failed.", error=str(exc))

    def save(self, path=None):
        service = self.ctx.state.get("word_service")
        if service is None:
            return DesktopActionResult("failed", self.application, "save", "No JARVIS Word document is active.")
        saved = service.save(path)
        return DesktopActionResult("success", self.application, "save", "Word document saved.", str(saved or path or ""))

    def export(self, format="pdf", path=None):
        if str(format).lower() != "pdf":
            return DesktopActionResult("failed", self.application, "export", "Word export supports PDF in this adapter.")
        service = self.ctx.state.get("word_service")
        if service is None:
            return DesktopActionResult("failed", self.application, "export", "No JARVIS Word document is active.")
        target = Path(path) if path else created_files_folder("PDF") / descriptive_filename("JARVIS Word Export", ".pdf")
        service.export_pdf(target)
        return DesktopActionResult("success", self.application, "export", "Word document exported to PDF.", str(target))


class ExcelAdapter(DesktopApplicationAdapter):
    application = "Microsoft Excel"

    def _service(self, create=False):
        service = self.ctx.state.get("excel_service")
        if service is None and create:
            from skills.office_service import ExcelService
            service = ExcelService()
            service.open(visible=True)
            self.ctx.state["excel_service"] = service
        return service

    def open_excel(self):
        return self._service(create=True).app

    def create_workbook(self):
        return self._service(create=True).new_workbook()

    def select_sheet(self, name):
        return self._service(create=True).select_sheet(name)

    def rename_sheet(self, current_name, new_name):
        return self._service(create=True).rename_sheet(current_name, new_name)

    def write_range(self, rows, start_row=1, start_col=1):
        return self._service(create=True).write_range(rows, start_row, start_col)

    def add_formula(self, cell, formula):
        return self._service(create=True).add_formula(cell, formula)

    def format_range(self, address, **kwargs):
        return self._service(create=True).format_range(address, **kwargs)

    def create_table(self, address, name="JarvisTable"):
        return self._service(create=True).create_table(address, name)

    def create_chart(self, address, chart_type=51):
        return self._service(create=True).create_chart(address, chart_type)

    def save_xlsx(self, path):
        return self._service(create=True).save(path)

    def export_csv(self, path):
        return self._service(create=True).export_csv(path)

    def close_workbook(self, save=True):
        service = self._service()
        return service.close_workbook(save) if service is not None else False

    def create_spreadsheet(self, topic="spreadsheet", mode="instant", save_after_completion=True):
        title = (topic or "Spreadsheet").strip().title()
        path = unique_path(created_files_folder("Excel") / descriptive_filename(title, ".xlsx"))
        task = getattr(self.ctx, "live_task", None)
        try:
            if task is not None:
                task.start(f"excel-{int(time.time())}", f"Create {topic}", application=self.application, mode=mode.upper())
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill

            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Monthly Budget" if "budget" in topic.lower() else title[:31]
            if "budget" in topic.lower():
                rows = [
                    ["Category", "Budget", "Actual", "Difference", "Notes"],
                    ["Housing", 0, 0, "=B2-C2", ""],
                    ["Utilities", 0, 0, "=B3-C3", ""],
                    ["Transport", 0, 0, "=B4-C4", ""],
                    ["Food", 0, 0, "=B5-C5", ""],
                    ["Total", "=SUM(B2:B5)", "=SUM(C2:C5)", "=B6-C6", ""],
                ]
            else:
                rows = [
                    ["Task ID", "Task", "Owner", "Priority", "Status", "Start Date", "Due Date", "Progress %", "Risk", "Next Action", "Notes"],
                    [1, title, "", "Medium", "Not Started", "", "", 0, "", "", ""],
                ]
            for row in rows:
                sheet.append(row)
            for cell in sheet[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="1F4E78")
            sheet.auto_filter.ref = sheet.dimensions
            sheet.freeze_panes = "A2"
            for column in sheet.columns:
                width = max(len(str(cell.value or "")) for cell in column) + 2
                sheet.column_dimensions[column[0].column_letter].width = min(width, 42)
            if save_after_completion:
                workbook.save(path)
            workbook.close()
            if save_after_completion and not path.exists():
                raise RuntimeError("Excel reported success but the XLSX file was not found")
            from skills.system_control import open_owned_file_in_application
            opened = open_owned_file_in_application(path, self.application, self.ctx)
            self.ctx.state["active_office_path"] = str(path)
            self.ctx.state["current_application"] = self.application
            if task is not None:
                task.complete()
            message = f"Created and verified {path.name}."
            message += " Excel remains open for review." if opened else " The file is ready, but Excel did not open in time."
            return DesktopActionResult("success", self.application, "create_spreadsheet", message, str(path))
        except Exception as exc:
            if task is not None:
                task.fail(exc)
            return DesktopActionResult("failed", self.application, "create_spreadsheet", "Spreadsheet creation failed.", error=str(exc))

    def save(self, path=None):
        service = self.ctx.state.get("excel_service")
        if service is None:
            return DesktopActionResult("failed", self.application, "save", "No JARVIS workbook is active.")
        saved = service.save(path)
        return DesktopActionResult("success", self.application, "save", "Workbook saved.", str(saved or path or ""))


class PowerPointAdapter(DesktopApplicationAdapter):
    application = "Microsoft PowerPoint"

    def _service(self, create=False):
        service = self.ctx.state.get("powerpoint_service")
        if service is None and create:
            from skills.office_service import PowerPointService
            service = PowerPointService()
            service.open(visible=True)
            self.ctx.state["powerpoint_service"] = service
        return service

    def open_powerpoint(self):
        return self._service(create=True).app

    def create_blank_presentation(self):
        return self._service(create=True).new_presentation()

    def add_slide(self, title, body=(), layout=2):
        return self._service(create=True).add_slide(title, body, layout)

    def set_slide_layout(self, slide, layout):
        slide.Layout = int(layout)
        return True

    def set_title(self, slide, title):
        slide.Shapes(1).TextFrame.TextRange.Text = str(title)
        return True

    def set_body(self, slide, body):
        slide.Shapes(2).TextFrame.TextRange.Text = str(body)
        return True

    def add_image(self, path, **kwargs):
        return self._service(create=True).add_image(path, **kwargs)

    def add_chart(self, *args, **kwargs):
        raise ValueError("PowerPoint charts must be created from a verified workbook data source")

    def add_notes(self, text, slide=None):
        return self._service(create=True).add_notes(text, slide)

    def run_slideshow(self):
        return self._service(create=True).run_slideshow()

    def save_pptx(self, path):
        return self._service(create=True).save(path)

    def export_pdf(self, path):
        return self._service(create=True).export_pdf(path)

    def close_presentation(self):
        service = self._service()
        return service.close_presentation() if service is not None else False

    def create_presentation(self, topic="Presentation", slides=10, mode="instant", save_after_completion=True):
        title = (topic or "Presentation").strip().title()
        path = unique_path(created_files_folder("PowerPoint") / descriptive_filename(f"{title} Presentation", ".pptx"))
        structure = (
            "Purpose", "Current Problem", "Key Facts", "Proposed Solution",
            "How It Works", "Benefits", "Risks", "Recommendations", "Final Message",
        )
        count = max(3, min(int(slides or 10), 20))
        try:
            from pptx import Presentation

            presentation = Presentation()
            title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            title_slide.shapes.title.text = title
            title_slide.placeholders[1].text = "Created by JARVIS"
            for heading in structure[:count - 1]:
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                slide.shapes.title.text = heading
                frame = slide.placeholders[1].text_frame
                frame.text = f"{heading} for {title}"
                for bullet in ("Key point", "Recommended next action"):
                    paragraph = frame.add_paragraph()
                    paragraph.text = bullet
                if mode in {"visible", "structured"}:
                    time.sleep(max(0, Config.LIVE_TYPING_DELAY_MS) / 1000.0)
            if save_after_completion:
                presentation.save(path)
            if save_after_completion and not path.exists():
                raise RuntimeError("PowerPoint reported success but the PPTX file was not found")
            from skills.system_control import open_owned_file_in_application
            opened = open_owned_file_in_application(path, self.application, self.ctx)
            self.ctx.state["active_office_path"] = str(path)
            self.ctx.state["current_application"] = self.application
            message = f"Created and verified {path.name}."
            message += " PowerPoint remains open for review." if opened else " The file is ready, but PowerPoint did not open in time."
            return DesktopActionResult("success", self.application, "create_presentation", message, str(path))
        except Exception as exc:
            return DesktopActionResult("failed", self.application, "create_presentation", "Presentation creation failed.", error=str(exc))

    def save(self, path=None):
        service = self.ctx.state.get("powerpoint_service")
        if service is None:
            return DesktopActionResult("failed", self.application, "save", "No JARVIS presentation is active.")
        saved = service.save(path)
        return DesktopActionResult("success", self.application, "save", "Presentation saved.", str(saved or path or ""))


class DesktopAdapterRegistry:
    def __init__(self, ctx):
        self._adapters = {
            "word": WordAdapter(ctx),
            "excel": ExcelAdapter(ctx),
            "powerpoint": PowerPointAdapter(ctx),
        }
        for key, application in {
            "outlook": "Microsoft Outlook", "onenote": "Microsoft OneNote",
            "access": "Microsoft Access", "teams": "Microsoft Teams",
            "paint": "Microsoft Paint", "edge": "Microsoft Edge",
        }.items():
            adapter = DesktopApplicationAdapter(ctx)
            adapter.application = application
            self._adapters[key] = adapter

    def get(self, application):
        low = str(application or "").lower().replace("microsoft", "").strip()
        if low in self._adapters:
            return self._adapters[low]
        raise ValueError(f"Unsupported desktop application adapter: {application}")

    def names(self):
        return tuple(sorted(self._adapters))


class DesktopAutomationService:
    def __init__(self, ctx):
        self.ctx = ctx
        self.adapters = DesktopAdapterRegistry(ctx)
        self.logger = DesktopActionLogger()

    def execute(self, intent):
        skill = intent.get("skill", "")
        params = dict(intent.get("params", {}) or {})
        intent_group = params.pop("intent_group", skill)
        # The command field is execution identity, not ambient UI state. A
        # pause/resume command may update ``last_command_text`` while this
        # slower action is still running, so bind the initiating command now.
        initiating_command = self.ctx.state.get("last_command_text", "")
        active_website = str(self.ctx.state.get("active_website", "")).lower()
        if skill == "office.create_document" and active_website == "google docs":
            adapter = self.ctx.website_adapters.get("google docs")
            document_type = params.get("document_type", "document")
            topic = params.get("topic", "") or document_type
            sections = self.adapters.get("word")._sections(document_type, topic)
            text = "\n\n".join(f"{heading}\n{body}" for heading, body in sections)
            adapter.create_blank(topic.title(), text)
            return f"Created the {document_type} visibly in Google Docs and left it open."
        if skill == "office.create_spreadsheet" and active_website == "google sheets":
            self.ctx.website_adapters.get("google sheets").create_monthly_budget()
            return "Created the monthly budget visibly in Google Sheets and left it open."
        if skill == "office.create_presentation" and active_website == "google slides":
            self.ctx.website_adapters.get("google slides").create_presentation(
                params.get("topic", "Presentation"), params.get("slides", 10)
            )
            return "Created the presentation visibly in Google Slides and left it open."
        result = None
        if skill == "office.create_document":
            result = self.adapters.get("word").create_document(**params)
        elif skill == "office.create_spreadsheet":
            result = self.adapters.get("excel").create_spreadsheet(**params)
        elif skill == "office.create_presentation":
            result = self.adapters.get("powerpoint").create_presentation(**params)
        elif skill in {"office.save", "office.export"}:
            service_key = next((key for key in ("word", "excel", "powerpoint") if self.ctx.state.get(f"{key}_service")), None)
            if service_key is None:
                result = DesktopActionResult("failed", "Microsoft Office", skill.split(".")[-1], "No JARVIS Office file is active.")
            else:
                adapter = self.adapters.get(service_key)
                operation = getattr(adapter, skill.split(".")[-1], None)
                if operation is None:
                    result = DesktopActionResult("failed", adapter.application, skill.split(".")[-1], "That export is not supported for the active application.")
                else:
                    result = operation(**{key: value for key, value in params.items() if key != "application"})
        else:
            raise ValueError(f"Unsupported desktop automation intent: {skill}")
        self.logger.write(intent=intent_group, application=result.application,
                          action=result.action, result=result.message,
                          error=result.error, file_path=result.file_path,
                          command=initiating_command)
        return result.message if result.status == "success" else f"{result.message} {result.error}".strip()
