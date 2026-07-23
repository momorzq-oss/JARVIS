from types import SimpleNamespace

import pytest

from config import Config
from core.live_task import LiveTaskController
from skills.desktop_automation import (
    DesktopActionResult, DesktopAdapterRegistry, DesktopAutomationService,
    ExcelAdapter, PowerPointAdapter, WordAdapter,
    created_files_folder,
    descriptive_filename,
    unique_path,
)


def _ctx():
    return SimpleNamespace(state={}, llm=SimpleNamespace(available=False))


def test_descriptive_filename_rejects_generic_and_invalid_characters():
    name = descriptive_filename("Staff: Training / Proposal", ".docx")
    assert name == "Staff Training Proposal.docx"
    assert name not in {"Document1.docx", "Book1.xlsx", "Presentation1.pptx"}


def test_created_files_use_required_documents_tree(tmp_path, monkeypatch):
    monkeypatch.setattr(Config, "CREATED_FILES_DIR", tmp_path / "Jarvis Created Files")
    folder = created_files_folder("Word")
    assert folder == tmp_path / "Jarvis Created Files" / "Word"
    assert folder.is_dir()


def test_unique_path_never_overwrites_existing_file(tmp_path):
    original = tmp_path / "Proposal.docx"
    original.write_text("existing", encoding="utf-8")
    assert unique_path(original) == tmp_path / "Proposal (2).docx"


def test_desktop_registry_exposes_dedicated_application_adapters():
    registry = DesktopAdapterRegistry(_ctx())
    names = registry.names()
    assert set(("word", "excel", "powerpoint", "outlook", "onenote", "access", "teams", "paint", "edge")) <= set(names)
    for method in (
        "open_word", "create_blank_document", "insert_text", "apply_heading",
        "insert_table", "save_docx", "export_pdf", "close_document",
    ):
        assert callable(getattr(registry.get("word"), method))
    for method in (
        "open_excel", "create_workbook", "select_sheet", "rename_sheet",
        "write_range", "add_formula", "format_range", "create_table",
        "create_chart", "save_xlsx", "export_csv", "close_workbook",
    ):
        assert callable(getattr(registry.get("excel"), method))
    for method in (
        "open_powerpoint", "create_blank_presentation", "add_slide",
        "set_slide_layout", "set_title", "set_body", "add_image",
        "add_notes", "run_slideshow", "save_pptx", "export_pdf",
        "close_presentation",
    ):
        assert callable(getattr(registry.get("powerpoint"), method))


def test_unknown_desktop_adapter_is_rejected():
    with pytest.raises(ValueError, match="Unsupported desktop application adapter"):
        DesktopAdapterRegistry(_ctx()).get("unknown")


def test_excel_creation_uses_deterministic_file_path_and_verified_launcher(tmp_path, monkeypatch):
    from openpyxl import load_workbook

    monkeypatch.setattr(Config, "CREATED_FILES_DIR", tmp_path)
    opened = []
    monkeypatch.setattr(
        "skills.system_control.open_owned_file_in_application",
        lambda path, application, _ctx: opened.append((path, application)) or True,
    )
    ctx = _ctx()
    result = ExcelAdapter(ctx).create_spreadsheet("monthly budget tracker")

    path = tmp_path / "Excel" / "Monthly Budget Tracker.xlsx"
    assert result.status == "success"
    assert opened == [(path, "Microsoft Excel")]
    workbook = load_workbook(path, data_only=False)
    assert workbook.active["D2"].value == "=B2-C2"
    assert workbook.active["B6"].value == "=SUM(B2:B5)"
    workbook.close()


def test_powerpoint_creation_uses_deterministic_file_path_and_verified_launcher(tmp_path, monkeypatch):
    from pptx import Presentation

    monkeypatch.setattr(Config, "CREATED_FILES_DIR", tmp_path)
    opened = []
    monkeypatch.setattr(
        "skills.system_control.open_owned_file_in_application",
        lambda path, application, _ctx: opened.append((path, application)) or True,
    )
    ctx = _ctx()
    result = PowerPointAdapter(ctx).create_presentation("local AI", slides=5)

    path = tmp_path / "PowerPoint" / "Local Ai Presentation.pptx"
    assert result.status == "success"
    assert opened == [(path, "Microsoft PowerPoint")]
    presentation = Presentation(path)
    assert len(presentation.slides) == 5
    assert presentation.slides[0].shapes.title.text == "Local Ai"


def test_desktop_action_log_keeps_initiating_command_during_resume():
    ctx = _ctx()
    ctx.state["last_command_text"] = "Create the original proposal."
    service = DesktopAutomationService(ctx)
    captured = []

    class Adapter:
        def create_document(self, **_params):
            ctx.state["last_command_text"] = "resume the current task"
            return DesktopActionResult(
                "success", "Microsoft Word", "create_document", "Created.",
            )

    service.adapters.get = lambda _name: Adapter()
    service.logger = SimpleNamespace(write=lambda **entry: captured.append(entry))

    assert service.execute({
        "skill": "office.create_document",
        "params": {"intent_group": "CREATE_DOCUMENT"},
    }) == "Created."
    assert captured[0]["command"] == "Create the original proposal."


def test_cancelled_word_creation_closes_unsaved_document_and_stays_cancelled(
        tmp_path, monkeypatch):
    monkeypatch.setattr(Config, "CREATED_FILES_DIR", tmp_path)
    closed = []
    task = LiveTaskController()
    ctx = SimpleNamespace(
        state={}, llm=SimpleNamespace(available=False), live_task=task,
    )

    class FakeWordService:
        process_id = 1
        window_handle = 2

        def open(self, visible=True):
            return visible

        def new_document(self):
            task.cancel()
            return "document"

        def insert_heading(self, *_args, **_kwargs):
            return True

        def close(self, save=False):
            closed.append(save)
            return True

    monkeypatch.setattr("skills.office_service.WordService", FakeWordService)

    result = WordAdapter(ctx).create_document(
        document_type="proposal", topic="cancellation validation",
        mode="structured",
    )

    assert result.status == "cancelled"
    assert task.snapshot()["status"] == "cancelled"
    assert closed == [False]
    assert not list(tmp_path.rglob("*.docx"))
